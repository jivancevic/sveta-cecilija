# -*- coding: utf-8 -*-
"""Build the dark cinematic board deck (16:9) from content.py."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from PIL import Image

import content as C

HERE = os.path.dirname(__file__)
ROOT = os.path.abspath(os.path.join(HERE, "..", "..", ".."))
ASSETS = os.path.join(ROOT, "docs", "presentation", "assets")
BG = os.path.join(ASSETS, "bg")
FR = os.path.join(ASSETS, "frames")
OUT = os.path.join(ROOT, "docs", "presentation", "out", "Digitalna-obnova-HGD-Sveta-Cecilija.pptx")

# ---- brand ----
DISPLAY = "Bodoni 72 Smallcaps"
BODY = "Helvetica Neue"
MONO = "Menlo"
GOLD = "B8881A"
GOLDB = "D9A526"
CREAM = "F4EFE6"
MUTE = "B7AD9B"
DARK = "0A0908"

EMU_IN = 914400
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


def slide_new():
    return prs.slides.add_slide(BLANK)


def add_bg(slide, name):
    slide.shapes.add_picture(os.path.join(BG, name), 0, 0, width=prs.slide_width, height=prs.slide_height)


def scrim(slide, x, y, w, h, hexc=DARK, alpha=30):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    sp.line.fill.background()
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor.from_string(hexc)
    srgb = sp._element.spPr.find(qn('a:solidFill')).find(qn('a:srgbClr'))
    srgb.append(srgb.makeelement(qn('a:alpha'), {'val': str(int(alpha * 1000))}))
    return sp


def rule(slide, x, y, w, hexc=GOLD, thick=0.022):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(thick))
    sp.shadow.inherit = False
    sp.line.fill.background()
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor.from_string(hexc)
    return sp


def _set_run(r, text, font, size, hexc, bold=False, italic=False, spc=None):
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = RGBColor.from_string(hexc)
    if spc is not None:
        r.font._rPr.set('spc', str(int(spc * 100)))


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    return tf


def para(tf, first=False, align=PP_ALIGN.LEFT, space_before=0, space_after=0, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    if line is not None:
        p.line_spacing = line
    return p


def add_pic_fit(slide, path, cx, cy, max_w, max_h):
    """Place an image centered in (cx,cy) box, scaled to fit max_w x max_h (inches)."""
    im = Image.open(path)
    iw, ih = im.size
    ar = iw / ih
    w = max_w
    h = w / ar
    if h > max_h:
        h = max_h
        w = h * ar
    left = cx - w / 2
    top = cy - h / 2
    slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(w), Inches(h))
    return left, top, w, h


def header(slide, eyebrow, title, x=0.92, y=0.82, w=11.5, title_size=46, title_color=GOLDB):
    tf = textbox(slide, x, y, w, 0.4)
    p = para(tf, first=True)
    _set_run(p.runs and p.add_run() or p.add_run(), eyebrow, MONO, 12.5, GOLD, spc=2.2)
    rule(slide, x, y + 0.34, 0.5)
    tf2 = textbox(slide, x, y + 0.5, w, 1.0)
    p2 = para(tf2, first=True, line=0.98)
    _set_run(p2.add_run(), title, DISPLAY, title_size, title_color)
    return y + 0.5 + (title_size / 72.0) * 1.25


def bullets(slide, points, x, y, w, lead_size=18, body_size=14.5, gap=12):
    """points: list of (lead, body) or strings."""
    tf = textbox(slide, x, y, w, SH - y - 0.6)
    first = True
    for it in points:
        if isinstance(it, (tuple, list)):
            lead, body = it
            p = para(tf, first=first, space_before=0 if first else gap, line=1.06)
            _set_run(p.add_run(), lead, BODY, lead_size, CREAM, bold=True)
            p2 = para(tf, line=1.12, space_before=2)
            _set_run(p2.add_run(), body, BODY, body_size, MUTE)
        else:
            p = para(tf, first=first, space_before=0 if first else gap, line=1.15)
            _set_run(p.add_run(), "›  ", BODY, body_size, GOLD, bold=True)
            _set_run(p.add_run(), it, BODY, body_size, CREAM)
        first = False
    return tf


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ---------------- layouts ----------------

def L_title(s, d):
    add_bg(s, d["bg"])
    scrim(s, 0, 0, SW, SH, alpha=22)
    tf = textbox(s, 0.95, 2.55, 11.4, 0.5)
    p = para(tf, first=True)
    _set_run(p.add_run(), d["eyebrow"], MONO, 13, GOLD, spc=2.4)
    rule(s, 0.97, 3.0, 0.6, GOLD)
    tf2 = textbox(s, 0.92, 3.15, 11.5, 1.5)
    p2 = para(tf2, first=True, line=0.95)
    _set_run(p2.add_run(), d["title"], DISPLAY, 80, GOLDB)
    tf3 = textbox(s, 0.95, 4.65, 10.4, 1.0)
    p3 = para(tf3, first=True, line=1.2)
    _set_run(p3.add_run(), d["subtitle"], BODY, 19, CREAM)
    tf4 = textbox(s, 0.95, 6.7, 11.4, 0.4)
    p4 = para(tf4, first=True)
    _set_run(p4.add_run(), d["footer"], MONO, 11.5, MUTE, spc=1.2)
    set_notes(s, d["notes"])


def L_statement(s, d):
    add_bg(s, d["bg"])
    scrim(s, 0, 0, SW, SH, alpha=28)
    tf = textbox(s, 0.95, 2.45, 11.4, 0.5)
    p = para(tf, first=True)
    _set_run(p.add_run(), d["eyebrow"], MONO, 12.5, GOLD, spc=2.4)
    rule(s, 0.97, 2.9, 0.6)
    tf2 = textbox(s, 0.92, 3.05, 11.4, 1.1)
    p2 = para(tf2, first=True, line=0.96)
    _set_run(p2.add_run(), d["title"], DISPLAY, 60, GOLDB)
    tf3 = textbox(s, 0.95, 4.35, 10.6, 1.6)
    p3 = para(tf3, first=True, line=1.3)
    _set_run(p3.add_run(), d["subtitle"], BODY, 20, CREAM)
    set_notes(s, d["notes"])


def L_divider(s, d):
    add_bg(s, d["bg"])
    scrim(s, 0, 0, SW, SH, alpha=30)
    tf = textbox(s, 0.95, 3.05, 11.4, 0.5)
    p = para(tf, first=True)
    _set_run(p.add_run(), d["kicker"], MONO, 13, GOLD, spc=3.0)
    rule(s, 0.97, 3.55, 0.7)
    tf2 = textbox(s, 0.92, 3.7, 11.4, 1.3)
    p2 = para(tf2, first=True, line=0.95)
    _set_run(p2.add_run(), d["title"], DISPLAY, 64, CREAM)
    set_notes(s, d["notes"] if "notes" in d else "")


def L_points_image(s, d):
    add_bg(s, d["bg"])
    scrim(s, 0, 0, SW, SH, alpha=20)
    header(s, d["eyebrow"], d["title"])
    bullets(s, d["points"], 0.92, 2.55, 6.1, lead_size=19, body_size=14, gap=15)
    # image on right
    l, t, w, h = add_pic_fit(s, os.path.join(FR, d["image"]), cx=10.05, cy=4.35, max_w=5.4, max_h=4.2)
    set_notes(s, d["notes"])


def L_compare(s, d):
    add_bg(s, d["bg"])
    scrim(s, 0, 0, SW, SH, alpha=30)
    header(s, d["eyebrow"], d["title"])
    # two frames side by side
    add_pic_fit(s, os.path.join(FR, d["left"]), cx=3.55, cy=4.35, max_w=5.5, max_h=3.7)
    add_pic_fit(s, os.path.join(FR, d["right"]), cx=9.78, cy=4.35, max_w=5.5, max_h=3.7)
    for cx, lab, col in [(3.55, d["left_label"], MUTE), (9.78, d["right_label"], GOLDB)]:
        tf = textbox(s, cx - 2.0, 6.55, 4.0, 0.5, align=PP_ALIGN.CENTER)
        p = para(tf, first=True, align=PP_ALIGN.CENTER)
        _set_run(p.add_run(), lab, MONO, 14, col, spc=3.0, bold=True)
    set_notes(s, d["notes"])


def L_showcase(s, d):
    add_bg(s, d["bg"])
    scrim(s, 0, 0, SW, SH, alpha=22)
    header(s, d["eyebrow"], d["title"])
    bullets(s, d["points"], 0.92, 2.7, 4.7, body_size=15, gap=14)
    add_pic_fit(s, os.path.join(FR, d["image"]), cx=9.3, cy=4.35, max_w=6.9, max_h=4.35)
    set_notes(s, d["notes"])


def L_showcase2(s, d):
    add_bg(s, d["bg"])
    scrim(s, 0, 0, SW, SH, alpha=22)
    header(s, d["eyebrow"], d["title"])
    add_pic_fit(s, os.path.join(FR, d["image_a"]), cx=3.35, cy=3.75, max_w=4.7, max_h=3.5)
    add_pic_fit(s, os.path.join(FR, d["image_b"]), cx=8.05, cy=3.75, max_w=4.7, max_h=3.5)
    # points across bottom
    tf = textbox(s, 0.92, 6.05, 11.5, 1.2)
    first = True
    for it in d["points"]:
        p = para(tf, first=first, space_before=0 if first else 6, line=1.1)
        _set_run(p.add_run(), "›  ", BODY, 14.5, GOLD, bold=True)
        _set_run(p.add_run(), it, BODY, 14.5, CREAM)
        first = False
    set_notes(s, d["notes"])


def L_qr(s, d):
    add_bg(s, d["bg"])
    scrim(s, 0, 0, SW, SH, alpha=24)
    header(s, d["eyebrow"], d["title"])
    # QR on a cream card at left
    qx, qy, qs = 1.15, 3.0, 3.0
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(qx - 0.18), Inches(qy - 0.18), Inches(qs + 0.36), Inches(qs + 0.36))
    card.shadow.inherit = False
    card.line.fill.background()
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
    s.shapes.add_picture(os.path.join(FR, d["image"]), Inches(qx), Inches(qy), Inches(qs), Inches(qs))
    # steps on right
    tf = textbox(s, 5.2, 2.9, 7.2, 3.6, anchor=MSO_ANCHOR.MIDDLE)
    first = True
    for num, body in d["steps"]:
        p = para(tf, first=first, space_before=0 if first else 16, line=1.1)
        _set_run(p.add_run(), num + ".  ", DISPLAY, 26, GOLDB, bold=True)
        _set_run(p.add_run(), body, BODY, 17, CREAM)
        first = False
    set_notes(s, d["notes"])


def L_channels(s, d):
    add_bg(s, d["bg"])
    scrim(s, 0, 0, SW, SH, alpha=22)
    header(s, d["eyebrow"], d["title"])
    # three channel cards
    n = len(d["channels"])
    cw, ch, gap = 3.5, 2.2, 0.45
    total = n * cw + (n - 1) * gap
    x0 = (SW - total) / 2
    cy = 2.75
    for i, (name, sub) in enumerate(d["channels"]):
        x = x0 + i * (cw + gap)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(cy), Inches(cw), Inches(ch))
        card.shadow.inherit = False
        card.line.color.rgb = RGBColor.from_string(GOLD)
        card.line.width = Pt(1.0)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor.from_string("141019")
        srgb = card._element.spPr.find(qn('a:solidFill')).find(qn('a:srgbClr'))
        srgb.append(srgb.makeelement(qn('a:alpha'), {'val': '55000'}))
        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _set_run(p.add_run(), name, DISPLAY, 30, GOLDB)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(6)
        for j, ln in enumerate(sub.split("\n")):
            if j:
                p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
            _set_run(p2.add_run(), ln, BODY, 13.5, MUTE)
    # points below
    tf = textbox(s, 0.92, 5.35, 11.5, 1.6)
    first = True
    for it in d["points"]:
        p = para(tf, first=first, space_before=0 if first else 8, line=1.12, align=PP_ALIGN.CENTER)
        _set_run(p.add_run(), it, BODY, 15, CREAM)
        first = False
    set_notes(s, d["notes"])


def L_points(s, d):
    add_bg(s, d["bg"])
    scrim(s, 0, 0, SW, SH, alpha=24)
    header(s, d["eyebrow"], d["title"])
    bullets(s, d["points"], 0.92, 2.85, 9.8, lead_size=21, body_size=16, gap=20)
    set_notes(s, d["notes"])


def L_recap(s, d):
    add_bg(s, d["bg"])
    scrim(s, 0, 0, SW, SH, alpha=34)
    header(s, d["eyebrow"], d["title"], title_size=52)
    # two-column checklist
    items = d["items"]
    half = (len(items) + 1) // 2
    cols = [items[:half], items[half:]]
    xs = [0.92, 6.95]
    for ci, col in enumerate(cols):
        tf = textbox(s, xs[ci], 2.95, 5.6, 4.0)
        first = True
        for it in col:
            p = para(tf, first=first, space_before=0 if first else 16, line=1.12)
            _set_run(p.add_run(), "✓  ", BODY, 17, GOLDB, bold=True)
            _set_run(p.add_run(), it, BODY, 17, CREAM)
            first = False
    set_notes(s, d["notes"])


def L_closing(s, d):
    add_bg(s, d["bg"])
    scrim(s, 0, 0, SW, SH, alpha=30)
    tf = textbox(s, 0.95, 2.7, 11.4, 0.5)
    p = para(tf, first=True)
    _set_run(p.add_run(), d["eyebrow"], MONO, 12.5, GOLD, spc=2.6)
    rule(s, 0.97, 3.15, 0.6)
    tf2 = textbox(s, 0.92, 3.3, 11.4, 1.2)
    p2 = para(tf2, first=True, line=0.96)
    _set_run(p2.add_run(), d["title"], DISPLAY, 66, GOLDB)
    tf3 = textbox(s, 0.95, 4.7, 10.6, 1.2)
    p3 = para(tf3, first=True, line=1.3)
    _set_run(p3.add_run(), d["subtitle"], BODY, 19, CREAM)
    tf4 = textbox(s, 0.95, 6.7, 11.4, 0.4)
    p4 = para(tf4, first=True)
    _set_run(p4.add_run(), d["footer"], MONO, 12, MUTE, spc=1.6)
    set_notes(s, d["notes"])


HANDLERS = {
    "title": L_title, "statement": L_statement, "divider": L_divider,
    "points-image": L_points_image, "compare": L_compare, "showcase": L_showcase,
    "showcase2": L_showcase2, "qr": L_qr, "channels": L_channels, "points": L_points,
    "recap": L_recap, "closing": L_closing,
}

for d in C.SLIDES:
    s = slide_new()
    HANDLERS[d["layout"]](s, d)

prs.save(OUT)
print("saved", OUT, "slides:", len(C.SLIDES))
